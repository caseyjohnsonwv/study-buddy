from abc import ABC, abstractmethod
import re
from typing import List
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.docstore.document import Document
from pptx import Presentation


class FileLoader(ABC):
    @abstractmethod
    def load_and_split(filename:str) -> List[Document]:
        pass


class FileLoaderPptx(FileLoader):

    def load_and_split(filename:str) -> List[Document]:
        assert filename.lower().endswith('.pptx')
        documents:List[Document] = []
        prs = Presentation(filename)
        for n,slide in enumerate(prs.slides):
            slide_texts = []
            # get all text from the slides themselves
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    slide_text_cleaned = re.sub('\s+', ' ', shape.text)
                    slide_texts.append(slide_text_cleaned)
            # also grab presenter's notes if available
            slide_notes = slide.notes_slide.notes_text_frame.text
            slide_notes_cleaned = re.sub('\s+', ' ', slide_notes)
            slide_texts.append(slide_notes_cleaned)
            full_slide_text = ' '.join(slide_texts)
            # convert slide to document for embedding
            metadata = {'filename':filename, 'filetype': '.pptx', 'slide_number':n+1}
            slide_document = Document(page_content=full_slide_text, metadata=metadata)
            documents.append(slide_document)
        # return final list of split documents
        text_splitter = RecursiveCharacterTextSplitter(
            separators=['\s+'],
            is_separator_regex=True,
            keep_separator=True,
            chunk_size=1000,
            chunk_overlap=200,
        )
        return text_splitter.split_documents(documents)
    