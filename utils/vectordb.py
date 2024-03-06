import glob
import os
import re
from typing import List, Set, Tuple
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.docstore.document import Document
from langchain_community.document_loaders.html_bs import BSHTMLLoader
from langchain_community.vectorstores.faiss import FAISS
from langchain_openai import OpenAIEmbeddings
from pptx import Presentation
import pypdf
import env

class VectorDB:
    _courses_root_path = None
    _embedding_function = None
    _vectorstore = None

    def search(self, phrase:str, k:int=5, metadata_filter:dict={}) -> List[Tuple[Document, float]]:
        return self._vectorstore.similarity_search_with_score(phrase, k=k, filter=metadata_filter)

    def __init__(self, courses_root_path:str):
        self._courses_root_path = os.path.normpath(courses_root_path)
        self._embedding_function = OpenAIEmbeddings(model='text-embedding-3-small', api_key=env.OPENAI_API_KEY)
        self._vectorstore = None
        # try to load from existing file
        if os.path.exists(os.path.join(self._courses_root_path, 'index.faiss')):
            print(f"Loading existing index from {self._courses_root_path}")
            self._vectorstore = FAISS.load_local(self._courses_root_path, self._embedding_function, index_name='index')
        # crawl course materials directory and add net new files
        print(f"Generating new embeddings from {self._courses_root_path}")
        filetype_loader_funcs = {
            '.pptx': self._load_pptx_docs,
            '.pdf': self._load_pdf_docs,
            '.html': self._load_html_docs,
        }
        final_docs_list = []
        for k,v in filetype_loader_funcs.items():
            if callable(v):
                print(f"Using function {v.__name__}() to index filetype {k}")
                glob_expr = os.path.normpath(os.path.join(self._courses_root_path, '*', '*', f"*{k}"))
                res_docs = v(glob_expr)
                print(f"Loaded {len(res_docs)} documents from {k} files")
                final_docs_list.extend(res_docs)
        if len(final_docs_list) > 0:
            # add documents to vectorstore if any were found
            if self._vectorstore is not None:
                self._vectorstore.add_documents(final_docs_list)
            # or create new one if needed
            else:
                self._vectorstore = FAISS.from_documents(final_docs_list, self._embedding_function)
        # export vectorstore to index file
        self._vectorstore.save_local(self._courses_root_path, index_name='index')

    def _get_indexed_files(self, filetype:str) -> Set[str]:
        if self._vectorstore is None:
            return set()
        retriever = self._vectorstore.as_retriever(search_kwargs={'k':999999})
        existing_documents = retriever.get_relevant_documents(query='', metadata={'filetype':filetype})
        return {os.path.join(doc.metadata['directory'], doc.metadata['filename']) for doc in existing_documents}
        
    def _load_pptx_docs(self, pptx_glob_expr:str) -> List[Document]:
        documents = []
        indexed_files = self._get_indexed_files(filetype='.pptx')
        # load documents from pptx files
        for f in glob.glob(pptx_glob_expr):
            if f not in indexed_files:
                prs = Presentation(f)
                filename, directory = os.path.basename(f), os.path.dirname(f)
                course_name = os.path.basename(os.path.dirname(directory))
                for n,slide in enumerate(prs.slides):
                    slide_texts = []
                    # get all text from the slides themselves
                    for shape in slide.shapes:
                        if hasattr(shape, 'text'):
                            slide_text_cleaned = re.sub('\s+', ' ', shape.text)
                            slide_texts.append(slide_text_cleaned)
                    # also grab presenter's notes if available
                    slide_notes = slide.notes_slide.notes_text_frame.text
                    slide_notes_cleaned = re.sub('\s+', ' ', slide_notes)
                    slide_texts.append(slide_notes_cleaned)
                    full_slide_text = ' '.join(slide_texts)
                    # convert slide to document for embedding
                    metadata = {'filename':filename, 'course_name': course_name, 'filetype': '.pptx', 'directory':directory, 'slide_number':n+1}
                    slide_document = Document(page_content=full_slide_text, metadata=metadata)
                    documents.append(slide_document)
        # split documents with text splitter
        text_splitter = RecursiveCharacterTextSplitter(separators=['\s+'], is_separator_regex=True, keep_separator=True, chunk_size=300, chunk_overlap=0)
        docs = text_splitter.split_documents(documents)
        # return list of documents ready for embedding
        return docs

    def _load_pdf_docs(self, pdf_glob_expr:str) -> List[Document]:
        documents = []
        indexed_files = self._get_indexed_files(filetype='.pdf')
        # load pdf files
        for f in glob.glob(pdf_glob_expr):
            if f not in indexed_files:
                filename, directory = os.path.basename(f), os.path.dirname(f)
                course_name = os.path.basename(os.path.dirname(directory))
                with open(f, 'rb') as pdf:
                    reader = pypdf.PdfReader(pdf)
                    for page in reader.pages:
                        page_text = page.extract_text()
                        page_text_cleaned = re.sub('\s+', ' ', page_text)
                        metadata = {'filename':filename, 'course_name': course_name, 'filetype':'.pdf', 'directory':directory, 'page_number':page.page_number}
                        page_document = Document(page_content=page_text_cleaned, metadata=metadata)
                        documents.append(page_document)
        # split documents with text splitter
        text_splitter = RecursiveCharacterTextSplitter(separators=['\s+'], is_separator_regex=True, keep_separator=True, chunk_size=500, chunk_overlap=0)
        docs = text_splitter.split_documents(documents)
        # return list of documents ready for embedding
        return docs
    
    def _load_html_docs(self, html_glob_expr:str) -> List[Document]:
        documents = []
        indexed_files = self._get_indexed_files(filetype='.html')
        # load html files
        for f in glob.glob(html_glob_expr):
            if f not in indexed_files:
                filename, directory = os.path.basename(f), os.path.dirname(f)
                course_name = os.path.basename(os.path.dirname(directory))
                metadata = {'filename':filename, 'course_name': course_name, 'filetype':'.html', 'directory':directory}
                loader = BSHTMLLoader(f, open_encoding='utf-8', bs_kwargs={'features':'html.parser'})
                html_docs = loader.load()
                for d in html_docs:
                    d.metadata = metadata
                documents.extend(html_docs)
        # split documents with text splitter
        text_splitter = RecursiveCharacterTextSplitter(separators=['\s+'], is_separator_regex=True, keep_separator=True, chunk_size=1000, chunk_overlap=0)
        docs = text_splitter.split_documents(documents)
        # return list of documents ready for embedding
        return docs
    